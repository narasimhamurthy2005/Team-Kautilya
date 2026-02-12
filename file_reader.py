import PyPDF2
import pytesseract
import cv2
import os
from docx import Document
from pptx import Presentation
from sumy.parsers.plaintext import PlaintextParser
from sumy.nlp.tokenizers import Tokenizer
from sumy.summarizers.lsa import LsaSummarizer

# --- AI Summarization Logic ---
def get_summary(text):
    """Generates a 2-sentence summary of the document content."""
    if not text or len(text) < 150:
        return text if text else "No readable content found."
    try:
        # Use LSA Summarizer for high-quality semantic extraction
        parser = PlaintextParser.from_string(text, Tokenizer("english"))
        summarizer = LsaSummarizer()
        summary = summarizer(parser.document, 2) 
        return " ".join([str(sentence) for sentence in summary])
    except Exception as e:
        return text[:200] + "..." # Fallback to snippet

# --- Enhanced Extraction Logic ---
def read_text_file(path):
    try:
        with open(path, 'r', encoding='utf-8') as f: return f.read()
    except: return ""

def read_pdf_file(path):
    text = ""
    try:
        reader = PyPDF2.PdfReader(path)
        for page in reader.pages:
            text += page.extract_text() or ""
    except: pass
    return text

def read_docx_file(path):
    try:
        doc = Document(path)
        return "\n".join([p.text for p in doc.paragraphs])
    except: return ""

def read_pptx_file(path):
    text = ""
    try:
        prs = Presentation(path)
        for s in prs.slides:
            for sh in s.shapes:
                if hasattr(sh, "text"): text += sh.text + " "
    except: pass
    return text

def read_image_file(path):
    """Winner Feature: Extracts text from PNG/JPG using OCR."""
    try:
        img = cv2.imread(path)
        # Convert to grayscale for better OCR accuracy
        gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
        return pytesseract.image_to_string(gray)
    except: return ""

def read_file(path):
    p = path.lower()
    if p.endswith(".txt"): return read_text_file(path)
    if p.endswith(".pdf"): return read_pdf_file(path)
    if p.endswith(".docx"): return read_docx_file(path)
    if p.endswith(".pptx"): return read_pptx_file(path)
    if p.endswith((".png", ".jpg", ".jpeg")): return read_image_file(path)
    return ""